"""① 拆解成 Deck(給 vision 用):
   - pptx 模式:抽文字/表格/講者備註 + 每張 slide 渲染成圖(需 pptx 工具鏈)
   - 圖片模式:直接把資料夾裡的投影片圖片當 slide(過渡方案,不需 pptx 工具鏈)

pptx 相依(python-pptx / pdf2image)與 Pillow 皆延遲載入,使圖片模式在未裝
pptx 工具鏈的環境仍可 import 執行。
"""
from __future__ import annotations

import base64
import io
import re
import shutil
import subprocess
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

from . import config

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}


@dataclass
class Slide:
    index: int                    # 1-based
    title: str = ""
    text: str = ""                # shape 內所有文字(含條列)
    tables: list[list[list[str]]] = field(default_factory=list)
    notes: str = ""               # 講者備註
    image_png: bytes = b""        # 整張 slide 的渲染圖

    def image_data_uri(self) -> str:
        b64 = base64.b64encode(self.image_png).decode()
        return f"data:image/png;base64,{b64}"


@dataclass
class Deck:
    path: Path
    slides: list[Slide]


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs if p.text).strip()


def _table(shape) -> list[list[str]] | None:
    if not shape.has_table:
        return None
    tbl = shape.table
    return [[cell.text.strip() for cell in row.cells] for row in tbl.rows]


def _render_slides_to_png(pptx_path: Path) -> list[bytes]:
    """soffice: pptx -> pdf -> 每頁 PNG。頁序對應 slide 序。"""
    from pdf2image import convert_from_path        # 延遲載入:圖片模式不需要

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise RuntimeError("找不到 soffice/libreoffice,請先安裝 LibreOffice。")
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp, str(pptx_path)],
            check=True, capture_output=True, timeout=600,
        )
        pdf = next(Path(tmp).glob("*.pdf"))
        images = convert_from_path(str(pdf), dpi=config.RENDER_DPI)
        out = []
        for img in images:
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            out.append(buf.getvalue())
        return out


def _natkey(name: str):
    """自然排序:數字段轉 int,使 slide_2 排在 slide_10 之前。"""
    return [int(s) if s.isdigit() else s.lower() for s in re.split(r"(\d+)", name)]


def _load_png(path: Path) -> bytes:
    """任意圖片 → PNG bytes(統一 mime,配合 Slide.image_data_uri)。Pillow 延遲載入。"""
    from PIL import Image

    with Image.open(path) as im:
        buf = io.BytesIO()
        im.convert("RGB").save(buf, format="PNG")
        return buf.getvalue()


def extract_image_dir(topic_dir: str | Path) -> Deck:
    """圖片模式:一個主題資料夾內的投影片圖片 → Deck(每圖一 slide,檔名自然排序)。"""
    topic_dir = Path(topic_dir)
    imgs = sorted(
        (p for p in topic_dir.iterdir() if p.is_file() and p.suffix.lower() in IMAGE_EXTS),
        key=lambda p: _natkey(p.name),
    )
    slides = [Slide(index=i, image_png=_load_png(p)) for i, p in enumerate(imgs, start=1)]
    return Deck(path=topic_dir, slides=slides)


def extract(pptx_path: str | Path) -> Deck:
    from pptx import Presentation                  # 延遲載入:圖片模式不需要

    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    pngs = _render_slides_to_png(pptx_path)

    slides: list[Slide] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts, tables, title = [], [], ""
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                title = _shape_text(shape)
                continue
            t = _shape_text(shape)
            if t:
                texts.append(t)
            tbl = _table(shape)
            if tbl:
                tables.append(tbl)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append(Slide(
            index=i, title=title, text="\n".join(texts), tables=tables, notes=notes,
            image_png=pngs[i - 1] if i - 1 < len(pngs) else b"",
        ))
    return Deck(path=pptx_path, slides=slides)
